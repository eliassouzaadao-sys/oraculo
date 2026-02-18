"""
Loaders para documentos: PDF, DOCX, XLSX, PPTX, TXT, CSV, JSON.
"""
import json
from typing import Optional
from langchain_community.document_loaders import PyPDFLoader, TextLoader, CSVLoader


def carrega_pdf(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo PDF.
    Tenta usar pdfplumber para melhor extracao de tabelas,
    com fallback para PyPDFLoader.

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto extraido do PDF
    """
    # Tenta usar pdfplumber para melhor extracao
    try:
        return _carrega_pdf_avancado(caminho)
    except ImportError:
        # Fallback para PyPDFLoader basico
        loader = PyPDFLoader(caminho)
        documentos = loader.load()
        return '\n\n'.join([doc.page_content for doc in documentos])


def _carrega_pdf_avancado(caminho: str) -> str:
    """
    Carrega PDF usando pdfplumber para preservar tabelas.
    Extrai texto normal e formata tabelas como Markdown.

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto extraido com tabelas formatadas
    """
    import pdfplumber

    partes = []

    with pdfplumber.open(caminho) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            partes_pagina = [f"[Pagina {page_num}]"]

            # Extrai texto normal
            texto = page.extract_text() or ""
            if texto.strip():
                partes_pagina.append(texto.strip())

            # Extrai e formata tabelas
            tabelas = page.extract_tables()
            for tab_idx, tabela in enumerate(tabelas):
                if tabela:
                    texto_tabela = _formata_tabela_markdown(tabela)
                    if texto_tabela:
                        partes_pagina.append(f"\n[Tabela {tab_idx + 1}]\n{texto_tabela}")

            if len(partes_pagina) > 1:
                partes.append('\n'.join(partes_pagina))

    return '\n\n'.join(partes)


def _formata_tabela_markdown(tabela: list) -> str:
    """
    Formata uma tabela como Markdown.

    Args:
        tabela: Lista de listas representando linhas da tabela

    Returns:
        Tabela formatada em Markdown
    """
    if not tabela or not tabela[0]:
        return ""

    linhas = []
    for i, row in enumerate(tabela):
        # Limpa celulas vazias ou None
        celulas = [str(c).strip() if c else "" for c in row]
        linha = "| " + " | ".join(celulas) + " |"
        linhas.append(linha)

        # Adiciona separador apos o header
        if i == 0:
            separador = "|" + "|".join(["---"] * len(celulas)) + "|"
            linhas.append(separador)

    return '\n'.join(linhas)


def carrega_txt(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo TXT.

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto do arquivo
    """
    loader = TextLoader(caminho, encoding='utf-8')
    documentos = loader.load()
    return '\n\n'.join([doc.page_content for doc in documentos])


def carrega_csv(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo CSV.

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto formatado do CSV
    """
    loader = CSVLoader(caminho, encoding='utf-8')
    documentos = loader.load()
    return '\n\n'.join([doc.page_content for doc in documentos])


def carrega_docx(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo DOCX (Microsoft Word).

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto extraido do documento
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError("Instale python-docx: pip install python-docx")

    doc = Document(caminho)

    partes = []

    # Extrai paragrafos
    for paragrafo in doc.paragraphs:
        texto = paragrafo.text.strip()
        if texto:
            partes.append(texto)

    # Extrai texto de tabelas
    for tabela in doc.tables:
        for linha in tabela.rows:
            celulas = [celula.text.strip() for celula in linha.cells if celula.text.strip()]
            if celulas:
                partes.append(" | ".join(celulas))

    return '\n\n'.join(partes)


def carrega_xlsx(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo XLSX (Microsoft Excel).

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto formatado das planilhas
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("Instale openpyxl: pip install openpyxl")

    workbook = openpyxl.load_workbook(caminho, data_only=True)

    partes = []

    for nome_planilha in workbook.sheetnames:
        planilha = workbook[nome_planilha]
        partes.append(f"=== Planilha: {nome_planilha} ===")

        for linha in planilha.iter_rows(values_only=True):
            celulas = [str(c) if c is not None else "" for c in linha]
            texto_linha = " | ".join(celulas).strip()
            if texto_linha and texto_linha != " | " * (len(celulas) - 1):
                partes.append(texto_linha)

    return '\n\n'.join(partes)


def carrega_pptx(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo PPTX (Microsoft PowerPoint).

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto extraido dos slides
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Instale python-pptx: pip install python-pptx")

    prs = Presentation(caminho)

    partes = []

    for i, slide in enumerate(prs.slides, 1):
        textos_slide = [f"=== Slide {i} ==="]

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                textos_slide.append(shape.text.strip())

            # Extrai texto de tabelas
            if shape.has_table:
                for linha in shape.table.rows:
                    celulas = [cell.text.strip() for cell in linha.cells if cell.text.strip()]
                    if celulas:
                        textos_slide.append(" | ".join(celulas))

        if len(textos_slide) > 1:
            partes.append('\n'.join(textos_slide))

    return '\n\n'.join(partes)


def carrega_json(caminho: str) -> str:
    """
    Carrega conteudo de um arquivo JSON.

    Args:
        caminho: Caminho do arquivo

    Returns:
        Texto formatado do JSON
    """
    with open(caminho, 'r', encoding='utf-8') as f:
        dados = json.load(f)

    return _json_para_texto(dados)


def _json_para_texto(dados, prefixo: str = "") -> str:
    """
    Converte estrutura JSON em texto legivel.

    Args:
        dados: Dados JSON
        prefixo: Prefixo para indentacao

    Returns:
        Texto formatado
    """
    partes = []

    if isinstance(dados, dict):
        for chave, valor in dados.items():
            if isinstance(valor, (dict, list)):
                partes.append(f"{prefixo}{chave}:")
                partes.append(_json_para_texto(valor, prefixo + "  "))
            else:
                partes.append(f"{prefixo}{chave}: {valor}")

    elif isinstance(dados, list):
        for i, item in enumerate(dados):
            if isinstance(item, (dict, list)):
                partes.append(f"{prefixo}[{i}]:")
                partes.append(_json_para_texto(item, prefixo + "  "))
            else:
                partes.append(f"{prefixo}- {item}")

    else:
        partes.append(f"{prefixo}{dados}")

    return '\n'.join(partes)


# Mapeamento de extensao para funcao de carregamento
LOADERS_DOCUMENTOS = {
    '.pdf': carrega_pdf,
    '.txt': carrega_txt,
    '.csv': carrega_csv,
    '.docx': carrega_docx,
    '.xlsx': carrega_xlsx,
    '.pptx': carrega_pptx,
    '.json': carrega_json,
}


def carrega_documento(caminho: str, extensao: Optional[str] = None) -> str:
    """
    Carrega documento baseado na extensao.

    Args:
        caminho: Caminho do arquivo
        extensao: Extensao do arquivo (opcional, detecta automaticamente)

    Returns:
        Texto extraido do documento

    Raises:
        ValueError: Se extensao nao suportada
    """
    if extensao is None:
        import os
        extensao = os.path.splitext(caminho)[1].lower()

    loader = LOADERS_DOCUMENTOS.get(extensao)
    if loader is None:
        raise ValueError(f"Extensao nao suportada: {extensao}")

    return loader(caminho)
